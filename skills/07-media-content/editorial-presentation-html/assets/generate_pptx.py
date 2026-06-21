"""
Editorial Presentation · PPTX Generator

把 Anthropic warm editorial 设计语言渲染到 .pptx。共享 design-tokens.md 配色 + typography.md
字体规则,通过 python-pptx 程序化生成 slides。图表先按数据关系选择组件,
不要把 KPI、时间窗、机制匹配和证据链默认画成同一种 proof bar。

Usage:
    from generate_pptx import EditorialDeck

    deck = EditorialDeck(
        title="血液科 IFI · v5",
        industry="medical",
        embed_fonts=False,  # True 嵌入 Fraunces+Inter+Mono;False 用 Cambria/Calibri/Consolas
    )
    deck.add_hero_slide(
        eyebrow="HEMATOLOGY · IFI · 2026",
        title_main="一句话起跑",
        title_accent="AI 编排 7 大血液疾病",
        sub="用户输入一句话疾病名 → 5 phase 流水线 → 4 件交付物",
        stats=[("38亿", "总市场"), ("7", "Modules"), ("22", "Charts"), ("43", "PMID")],
    )
    deck.add_evidence_slide(
        eyebrow="CITE-OR-BLOCK · 语义层",
        title="每一个数字 · 都能点击回溯",
        funnel=[("PubMed", "数千"), ("初筛", "~500"), ("入库", "47"), ("纳入 A/B/C", "43")],
        pyramid=[("最高级", 18, "green"), ("重要级", 17, "blue"), ("参考级", 7, "gold"), ("国际", 2, "neutral"), ("估算", "按需", "deep")],
    )
    deck.save("output/v5.pptx")

要求:
    pip install python-pptx Pillow

源页面参考:C:\\Users\\qiyon\\Desktop\\血液科市场调研_v5_desktop.html
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree

# ============================================================
# DESIGN TOKENS · 与 design-tokens.md :root 完全对应
# ============================================================

COLORS = {
    "bg_0":          RGBColor(250, 249, 245),  # #FAF9F5
    "bg_1":          RGBColor(243, 238, 227),  # #F3EEE3
    "bg_2":          RGBColor(235, 229, 216),  # #EBE5D8
    "bg_surface":    RGBColor(255, 255, 255),  # #FFFFFF
    "text_0":        RGBColor(20, 20, 19),
    "text_1":        RGBColor(44, 42, 38),
    "text_2":        RGBColor(91, 84, 74),
    "text_3":        RGBColor(98, 91, 79),
    "accent":        RGBColor(204, 120, 92),   # #CC785C
    "accent_hot":    RGBColor(232, 168, 136),
    "accent_deep":   RGBColor(164, 89, 70),
    "brand_blue":    RGBColor(46, 92, 138),
    "brand_purple":  RGBColor(107, 78, 143),
    "brand_gold":    RGBColor(184, 144, 60),
    "brand_green":   RGBColor(92, 141, 92),
    "brand_red":     RGBColor(194, 89, 74),
    "brand_pink":    RGBColor(199, 122, 154),
    "border_soft":   RGBColor(227, 220, 204),
    "border":        RGBColor(208, 200, 181),
    "white":         RGBColor(255, 255, 255),
}

INDUSTRY_VERTICAL = {
    "medical":   RGBColor(140, 43, 58),    # #8C2B3A 深红血色
    "tech":      RGBColor(26, 58, 110),    # #1A3A6E 深海军蓝
    "finance":   RGBColor(45, 90, 58),     # #2D5A3A 深森林绿
    "education": RGBColor(74, 45, 92),     # #4A2D5C 深紫罗兰
    "fashion":   RGBColor(122, 41, 64),    # #7A2940 深酒红
}

FONTS_EMBED = {
    "serif": "Fraunces",
    "sans":  "Inter",
    "mono":  "JetBrains Mono",
}
FONTS_FALLBACK = {
    "serif": "Cambria",
    "sans":  "Calibri",
    "mono":  "Consolas",
}

# 16:9 1920×1080 in EMU(1 inch = 914400 EMU,1 cm = 360000 EMU)
SLIDE_W_16_9 = Inches(13.333)  # 1920 px @ 144 dpi
SLIDE_H_16_9 = Inches(7.5)
SLIDE_W_4_3  = Inches(10)
SLIDE_H_4_3  = Inches(7.5)


# ============================================================
# EDITORIAL DECK CLASS
# ============================================================

class EditorialDeck:
    """Anthropic warm editorial 风格 PPTX 生成器。"""

    def __init__(
        self,
        title="Editorial Presentation",
        industry="medical",
        embed_fonts=False,
        aspect_ratio="16:9",
    ):
        self.title = title
        self.industry = industry
        self.fonts = FONTS_EMBED if embed_fonts else FONTS_FALLBACK
        self.colors = dict(COLORS)
        self.colors["industry_vertical"] = INDUSTRY_VERTICAL.get(industry, INDUSTRY_VERTICAL["medical"])

        self.prs = Presentation()
        if aspect_ratio == "16:9":
            self.prs.slide_width = SLIDE_W_16_9
            self.prs.slide_height = SLIDE_H_16_9
        else:
            self.prs.slide_width = SLIDE_W_4_3
            self.prs.slide_height = SLIDE_H_4_3

        self._blank_layout = self.prs.slide_layouts[6]  # blank layout

    # ====== 公共 slide 添加方法 ======

    def add_hero_slide(self, eyebrow, title_main, title_accent, sub, stats=None, jump_links=None):
        """Slide HERO · 大标题 + stat-bar"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, "bg_0")
        self._add_radial_glow(slide)

        # eyebrow(顶部胶囊)
        self._add_glass_eyebrow(slide, x=0.7, y=0.6, text=eyebrow)

        # hero title — 主行 + accent 行(后者用渐变文字)
        self._add_text(
            slide, x=0.7, y=1.2, w=12, h=1.1,
            text=title_main,
            font=self.fonts["serif"], size=54, bold=True,
            color=self.colors["text_0"], spacing_pt=-2.0,
        )
        self._add_gradient_text(
            slide, x=0.7, y=2.3, w=12, h=1.1,
            text=title_accent,
            font=self.fonts["serif"], size=54, bold=True,
        )

        # subtitle
        self._add_text(
            slide, x=0.7, y=3.6, w=11, h=0.9,
            text=sub,
            font=self.fonts["sans"], size=14,
            color=self.colors["text_2"], line_spacing=1.5,
        )

        # stat-bar
        if stats:
            self._draw_stat_bar(slide, x=0.7, y=5.0, w=12, h=1.5, stats=stats)

        # jump links
        if jump_links:
            self._add_jump_row(slide, x=0.7, y=6.7, links=jump_links)

        return slide

    def add_why_slide(self, eyebrow, title, title_accent, bar_chart=None, data_pain=None):
        """Slide WHY · 双柱图叙事"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, "bg_1")
        self._add_radial_glow(slide, color=self.colors["industry_vertical"], opacity=0.08)

        self._add_eyebrow_tag(slide, x=0.7, y=0.6, text=eyebrow, color="industry_vertical")

        # 标题(主白 + accent 用 industry vertical)
        self._add_split_title(slide, x=0.7, y=1.2, main=title, accent=title_accent, accent_color="industry_vertical")

        # 双柱图
        if bar_chart:
            self._draw_double_bar_chart(slide, x=0.7, y=2.5, w=12, h=2.8, data=bar_chart)

        # data-pain callout
        if data_pain:
            self._draw_data_pain_callout(slide, x=0.7, y=5.6, w=12, h=0.8, text=data_pain)

        return slide

    def add_evidence_slide(self, eyebrow, title, funnel=None, pyramid=None):
        """Slide EVIDENCE · 漏斗 + 金字塔并排"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, "bg_1")

        self._add_eyebrow_tag(slide, x=0.7, y=0.6, text=eyebrow)
        self._add_text(
            slide, x=0.7, y=1.2, w=12, h=0.9,
            text=title,
            font=self.fonts["serif"], size=42, bold=True,
            color=self.colors["text_0"], spacing_pt=-1.2,
        )

        # 双栏:漏斗左 / 金字塔右
        if funnel:
            self._draw_funnel(slide, x=0.7, y=2.6, w=5.8, h=4.2, stages=funnel)
        if pyramid:
            self._draw_evidence_pyramid(slide, x=7.0, y=2.6, w=5.8, h=4.2, tiers=pyramid)

        return slide

    def add_architecture_slide(self, eyebrow, title, phases=None, deliverables=None, iron_rule=None):
        """Slide ARCHITECTURE · phase pill 链 + delivery box"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, "bg_0")

        self._add_eyebrow_tag(slide, x=0.7, y=0.6, text=eyebrow)
        self._add_text(
            slide, x=0.7, y=1.2, w=12, h=0.9,
            text=title,
            font=self.fonts["serif"], size=42, bold=True,
            color=self.colors["text_0"], spacing_pt=-1.2,
        )

        if phases:
            self._draw_phase_pill_row(slide, x=0.7, y=2.6, w=12, h=1.4, phases=phases)
        if deliverables:
            self._draw_delivery_box(slide, x=0.7, y=4.4, w=12, h=1.3, items=deliverables)
        if iron_rule:
            self._draw_iron_rule_band(slide, x=0.7, y=6.1, w=12, h=0.6, text=iron_rule)

        return slide

    def add_proof_slide(self, eyebrow, title, proof_bars=None, key_message=None):
        """Slide PROOF · 同一指标横向比较专用的校验条"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, "bg_1")

        self._add_eyebrow_tag(slide, x=0.7, y=0.6, text=eyebrow)
        self._add_text(
            slide, x=0.7, y=1.2, w=12, h=0.9,
            text=title,
            font=self.fonts["serif"], size=42, bold=True,
            color=self.colors["text_0"], spacing_pt=-1.2,
        )

        if proof_bars:
            for i, bar in enumerate(proof_bars):
                self._draw_proof_bar(
                    slide, x=0.7, y=2.6 + i * 0.7, w=12, h=0.55,
                    name=bar["name"], domain=bar.get("domain", ""),
                    pct=bar["pct"], status=bar["status"],
                )

        if key_message:
            self._draw_data_pain_callout(slide, x=0.7, y=6.5, w=12, h=0.7, text=key_message, color="brand_green")

        return slide

    def add_relationship_slide(self, eyebrow, title, timeline=None, matrix=None, evidence=None, note=None):
        """Slide RELATIONSHIPS · 时间窗 + 机制矩阵 + 证据列表(非 bar 默认)。"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, "bg_0")
        self._add_radial_glow(slide, color=self.colors["industry_vertical"], opacity=0.06)

        self._add_eyebrow_tag(slide, x=0.7, y=0.6, text=eyebrow, color="industry_vertical")
        self._add_text(
            slide, x=0.7, y=1.2, w=12, h=0.9,
            text=title,
            font=self.fonts["serif"], size=42, bold=True,
            color=self.colors["text_0"], spacing_pt=-1.2,
        )

        if note:
            self._draw_data_pain_callout(slide, x=0.7, y=2.2, w=12, h=0.55, text=note, color="accent")
            top_y = 3.0
        else:
            top_y = 2.55

        if timeline:
            self._draw_decision_timeline(slide, x=0.7, y=top_y, w=12, h=1.45, items=timeline)

        lower_y = top_y + 1.75
        if matrix:
            self._draw_matrix(slide, x=0.7, y=lower_y, w=7.7, h=2.55, rows=matrix)
        if evidence:
            self._draw_evidence_list(slide, x=8.7, y=lower_y, w=4.0, h=2.55, items=evidence)

        return slide

    def add_cta_slide(self, eyebrow, team=None, tech_stack=None, ctas=None, contact=None):
        """Slide CTA · 团队 + 3 CTA cards + 联系"""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, "bg_2")
        self._add_radial_glow(slide, position="bottom")

        if eyebrow:
            self._add_eyebrow_tag(slide, x=4, y=0.5, text=eyebrow)

        if team:
            team_text = " · ".join(team)
            self._add_gradient_text(
                slide, x=0.7, y=1.0, w=12, h=1.6,
                text=team_text,
                font=self.fonts["serif"], size=36, bold=True,
                align="center",
            )

        if tech_stack:
            self._draw_tech_stack_row(slide, x=0.7, y=2.7, w=12, h=0.4, items=tech_stack)

        if ctas:
            cols = len(ctas)
            col_w = 12 / cols
            for i, cta in enumerate(ctas):
                self._draw_cta_card(
                    slide, x=0.7 + i * col_w, y=3.4, w=col_w - 0.2, h=2.6,
                    tag=cta["tag"], time=cta["time"], title=cta["title"], desc=cta["desc"],
                    primary=cta.get("primary", i < 2),
                )

        if contact:
            self._add_text(
                slide, x=0.7, y=6.5, w=12, h=0.4,
                text=f"问题或 demo 请求 · {contact}",
                font=self.fonts["sans"], size=10,
                color=self.colors["text_3"], align="center",
            )

        return slide

    def save(self, path):
        """导出 .pptx"""
        self.prs.save(path)
        return path

    # ====== 视觉组件绘制(私有方法)======

    def _fill_background(self, slide, color_key="bg_0"):
        """slide 填充底色"""
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = self.colors[color_key]

    def _add_radial_glow(self, slide, color=None, position="top-right", opacity=0.10):
        """添加径向渐变光晕(用透明椭圆 shape 模拟)"""
        if color is None:
            color = self.colors["accent"]
        positions = {
            "top-right":    (Inches(8), Inches(-2),  Inches(8), Inches(6)),
            "bottom-left":  (Inches(-2), Inches(4),  Inches(8), Inches(6)),
            "bottom":       (Inches(2.5), Inches(5), Inches(8), Inches(5)),
        }
        x, y, w, h = positions.get(position, positions["top-right"])
        shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        # 设置透明
        sp = shape.fill.fore_color._xFill
        alpha_ele = etree.SubElement(sp, qn("a:alpha"), {"val": str(int(opacity * 100000))})

    def _add_glass_eyebrow(self, slide, x, y, text):
        """玻璃态 eyebrow(顶部胶囊形)"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(4), Inches(0.4)
        )
        shape.adjustments[0] = 0.5  # 高圆角(胶囊)
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors["white"]
        shape.line.color.rgb = self.colors["border_soft"]
        shape.line.width = Pt(0.75)

        tf = shape.text_frame
        tf.margin_left = Inches(0.15)
        tf.margin_right = Inches(0.15)
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = self.fonts["sans"]
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = self.colors["accent_deep"]

    def _add_eyebrow_tag(self, slide, x, y, text, color="accent_deep"):
        """普通 section-tag(描边小胶囊)"""
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(3.5), Inches(0.35)
        )
        shape.adjustments[0] = 0.5
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors["bg_surface"]
        shape.line.color.rgb = self.colors["border"]
        shape.line.width = Pt(0.75)

        tf = shape.text_frame
        tf.margin_left = Inches(0.12)
        tf.margin_right = Inches(0.12)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.name = self.fonts["sans"]
        run.font.size = Pt(8)
        run.font.bold = True
        run.font.color.rgb = self.colors[color] if color in self.colors else self.colors["accent_deep"]

    def _add_text(self, slide, x, y, w, h, text, font, size, color=None,
                  bold=False, italic=False, align="left", spacing_pt=0, line_spacing=1.2):
        """通用文字框"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        if color:
            run.font.color.rgb = color
        if spacing_pt:
            run.font.spacing = Pt(spacing_pt)
        return tb

    def _add_gradient_text(self, slide, x, y, w, h, text, font, size, bold=True, align="left"):
        """渐变文字(深蓝→赭石→金 三色 stop)"""
        tb = self._add_text(slide, x, y, w, h, text, font, size, bold=bold, align=align)
        # 写入 XML 添加 gradient fill
        tf = tb.text_frame
        run = tf.paragraphs[0].runs[0]
        rPr = run._r.get_or_add_rPr()
        # 移除已有 fill
        for sf in rPr.findall(qn("a:solidFill")):
            rPr.remove(sf)
        # 添加 gradient
        gradFill = etree.SubElement(rPr, qn("a:gradFill"))
        gsLst = etree.SubElement(gradFill, qn("a:gsLst"))
        # 深蓝 0%
        gs1 = etree.SubElement(gsLst, qn("a:gs"), {"pos": "0"})
        clr1 = etree.SubElement(gs1, qn("a:srgbClr"), {"val": "1E3A5F"})
        # 赭石 55%
        gs2 = etree.SubElement(gsLst, qn("a:gs"), {"pos": "55000"})
        clr2 = etree.SubElement(gs2, qn("a:srgbClr"), {"val": "CC785C"})
        # 金 100%
        gs3 = etree.SubElement(gsLst, qn("a:gs"), {"pos": "100000"})
        clr3 = etree.SubElement(gs3, qn("a:srgbClr"), {"val": "B8903C"})
        # 线性方向 135 度
        lin = etree.SubElement(gradFill, qn("a:lin"), {"ang": "8100000", "scaled": "1"})
        return tb

    def _add_split_title(self, slide, x, y, main, accent, accent_color="accent"):
        """标题:主部分黑色 + accent 部分用色"""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(12), Inches(1.1))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run1 = p.add_run()
        run1.text = main
        run1.font.name = self.fonts["serif"]
        run1.font.size = Pt(42)
        run1.font.bold = True
        run1.font.color.rgb = self.colors["text_0"]
        run1.font.spacing = Pt(-1.2)

        run2 = p.add_run()
        run2.text = accent
        run2.font.name = self.fonts["serif"]
        run2.font.size = Pt(42)
        run2.font.bold = True
        run2.font.color.rgb = self.colors[accent_color] if accent_color in self.colors else self.colors["accent"]
        run2.font.spacing = Pt(-1.2)

    def _draw_stat_bar(self, slide, x, y, w, h, stats):
        """大数字 stat-bar(顶/底 1px 横线 + 4-5 个数字)"""
        n = len(stats)
        col_w = w / n
        # 顶线
        self._add_horizontal_line(slide, x, y, w, self.colors["border_soft"])
        # 底线
        self._add_horizontal_line(slide, x, y + h, w, self.colors["border_soft"])

        for i, (num, label, *caption) in enumerate([(s[0], s[1], s[2] if len(s) > 2 else None) for s in stats]):
            cx = x + i * col_w
            # 大数字
            self._add_text(
                slide, x=cx + 0.1, y=y + 0.2, w=col_w - 0.2, h=0.8,
                text=num,
                font=self.fonts["serif"], size=36, bold=True,
                color=self.colors["text_0"], spacing_pt=-1.5,
            )
            # label
            self._add_text(
                slide, x=cx + 0.1, y=y + 1.0, w=col_w - 0.2, h=0.3,
                text=label,
                font=self.fonts["sans"], size=8, bold=True,
                color=self.colors["text_2"], spacing_pt=0.6,
            )
            if caption[0]:
                self._add_text(
                    slide, x=cx + 0.1, y=y + 1.3, w=col_w - 0.2, h=0.2,
                    text=caption[0],
                    font=self.fonts["sans"], size=8,
                    color=self.colors["text_2"],
                )

    def _draw_double_bar_chart(self, slide, x, y, w, h, data):
        """
        双柱图 · data 格式:
            {
                "axis": ["阶段→", "阶段1", "阶段2", ...],  # N+1 列(第一列是 label)
                "rows": [
                    {"label": "Line A", "pct": "~50%", "cells": [{"type":"s-prevent","text":"预防"}, ...]},
                    ...
                ]
            }
        """
        # 卡片背景
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.adjustments[0] = 0.04
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors["bg_surface"]
        card.line.color.rgb = self.colors["border_soft"]

        axis = data["axis"]
        rows = data["rows"]
        n_cols = len(axis) - 1
        col_label_w = 1.2
        chart_w = w - 0.4 - col_label_w
        col_w = chart_w / n_cols
        inner_x = x + 0.3 + col_label_w
        inner_y = y + 0.3
        # axis labels
        for i, label in enumerate(axis[1:]):
            self._add_text(
                slide, x=inner_x + i * col_w, y=inner_y, w=col_w - 0.05, h=0.3,
                text=label,
                font=self.fonts["sans"], size=8, bold=True,
                color=self.colors["text_3"], align="center", spacing_pt=0.5,
            )
        # rows
        row_h = (h - 0.7) / len(rows)
        for ri, row in enumerate(rows):
            row_y = inner_y + 0.4 + ri * row_h
            # 左侧 meta
            self._add_text(
                slide, x=x + 0.2, y=row_y, w=col_label_w, h=0.4,
                text=row["pct"],
                font=self.fonts["serif"], size=18, bold=True,
                color=self.colors["text_0"], align="right",
            )
            self._add_text(
                slide, x=x + 0.2, y=row_y + 0.4, w=col_label_w, h=0.3,
                text=row["label"],
                font=self.fonts["sans"], size=7,
                color=self.colors["text_3"], align="right",
            )
            # cells
            for ci, cell in enumerate(row["cells"]):
                self._draw_bar_cell(
                    slide,
                    x=inner_x + ci * col_w,
                    y=row_y,
                    w=col_w - 0.05,
                    h=row_h - 0.1,
                    cell=cell,
                )

    def _draw_bar_cell(self, slide, x, y, w, h, cell):
        """双柱图单元格"""
        cell_type = cell.get("type", "default") if isinstance(cell, dict) else cell
        text = cell.get("text", "") if isinstance(cell, dict) else ""
        chips = cell.get("chips", []) if isinstance(cell, dict) else []

        type_to_color = {
            "s-prevent":  ("brand_blue",   0.18),
            "s-empiric":  ("brand_gold",   0.16),
            "s-dd":       ("brand_gold",   0.20),
            "s-target":   ("accent",       0.20),
            "s-maint":    ("brand_green",  0.16),
            "event":      ("brand_red",    0.15),
            "empty":      None,
            "default":    ("bg_1",         1.0),
        }
        color_info = type_to_color.get(cell_type, type_to_color["default"])

        if color_info is None:
            # empty cell:dashed border 条纹
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x), Inches(y), Inches(w), Inches(h))
            shape.adjustments[0] = 0.08
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors["bg_1"]
            shape.line.color.rgb = self.colors["border"]
            shape.line.width = Pt(0.75)
            # 透明度
            sp = shape.fill.fore_color._xFill
            etree.SubElement(sp, qn("a:alpha"), {"val": "35000"})
            return

        color_key, alpha = color_info
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.adjustments[0] = 0.08
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors[color_key]
        sp = shape.fill.fore_color._xFill
        etree.SubElement(sp, qn("a:alpha"), {"val": str(int(alpha * 100000))})
        shape.line.color.rgb = self.colors[color_key]
        shape.line.width = Pt(0.75)

        # 文字
        tf = shape.text_frame
        tf.margin_left = Inches(0.06)
        tf.margin_right = Inches(0.06)
        tf.margin_top = Inches(0.06)
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = self.fonts["sans"]
        run.font.size = Pt(9)
        run.font.bold = True
        run.font.color.rgb = self.colors["text_0"]

    def _draw_evidence_pyramid(self, slide, x, y, w, h, tiers):
        """5 级金字塔 · 从顶到底宽度递增"""
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors["bg_surface"]
        card.line.color.rgb = self.colors["border_soft"]

        # 标题
        self._add_text(
            slide, x=x + 0.3, y=y + 0.2, w=w - 0.6, h=0.4,
            text="🏆 证据金字塔",
            font=self.fonts["serif"], size=14, bold=True,
            color=self.colors["brand_green"],
        )

        widths = [0.6, 0.72, 0.84, 0.92, 1.0]
        color_map = {"green": "brand_green", "blue": "brand_blue", "gold": "accent",
                     "neutral": "bg_2", "deep": "accent_deep"}
        row_h = (h - 0.8) / max(len(tiers), 1)
        for i, tier in enumerate(tiers):
            label, count, key = tier[0], tier[1], tier[2] if len(tier) > 2 else "green"
            row_w = w * widths[i] if i < len(widths) else w
            row_x = x + (w - row_w) / 2
            row_y = y + 0.7 + i * row_h
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(row_x), Inches(row_y), Inches(row_w), Inches(row_h - 0.05))
            shape.adjustments[0] = 0.1
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors[color_map.get(key, "brand_green")]
            shape.line.fill.background()

            tf = shape.text_frame
            tf.margin_left = Inches(0.15)
            tf.margin_right = Inches(0.15)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"{label}    {count}"
            run.font.name = self.fonts["sans"]
            run.font.size = Pt(11)
            run.font.bold = True
            run.font.color.rgb = self.colors["white"] if key != "neutral" else self.colors["text_1"]

    def _draw_funnel(self, slide, x, y, w, h, stages):
        """文献漏斗 · 5 行从宽到窄"""
        # 卡片背景
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors["bg_surface"]
        card.line.color.rgb = self.colors["border_soft"]

        # 标题
        self._add_text(
            slide, x=x + 0.3, y=y + 0.2, w=w - 0.6, h=0.4,
            text="📚 检索漏斗",
            font=self.fonts["serif"], size=14, bold=True,
            color=self.colors["accent_deep"],
        )

        widths = [1.0, 0.95, 0.85, 0.75, 0.65]
        bg_colors = ["#F7EFE7", "#F2E2D3", "#ECCBB2", "#D99975", "#A45946"]
        text_colors = ["text_0", "text_0", "text_0", "white", "white"]
        row_h = (h - 0.8) / max(len(stages), 1)
        for i, (label, count) in enumerate(stages):
            row_w = w * widths[i] if i < len(widths) else w * 0.65
            row_x = x + (w - row_w) / 2
            row_y = y + 0.7 + i * row_h

            shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(row_x), Inches(row_y), Inches(row_w), Inches(row_h - 0.04))
            shape.fill.solid()
            hex_color = bg_colors[i] if i < len(bg_colors) else bg_colors[-1]
            shape.fill.fore_color.rgb = RGBColor.from_string(hex_color.lstrip('#'))
            shape.line.fill.background()

            tf = shape.text_frame
            tf.margin_left = Inches(0.2)
            tf.margin_right = Inches(0.2)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            run1 = p.add_run()
            run1.text = label + "    "
            run1.font.name = self.fonts["sans"]
            run1.font.size = Pt(10)
            run1.font.bold = True
            color_key = text_colors[i] if i < len(text_colors) else "white"
            run1.font.color.rgb = self.colors[color_key]
            run2 = p.add_run()
            run2.text = str(count)
            run2.font.name = self.fonts["serif"]
            run2.font.size = Pt(13)
            run2.font.bold = True
            run2.font.color.rgb = self.colors[color_key]

    def _draw_phase_pill_row(self, slide, x, y, w, h, phases):
        """5-6 个 phase pill 横向链"""
        n = len(phases)
        col_w = w / n
        for i, phase in enumerate(phases):
            ph_id, ph_name, ph_zh = phase[0], phase[1], phase[2] if len(phase) > 2 else ""
            cx = x + i * col_w + 0.1
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cx), Inches(y), Inches(col_w - 0.2), Inches(h))
            shape.adjustments[0] = 0.1
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors["bg_surface"]
            shape.line.color.rgb = self.colors["border"]
            self._add_card_shadow(shape)

            tf = shape.text_frame
            tf.margin_left = Inches(0.1)
            tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.1)

            # PHASE ID(紫色 mono 小字)
            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run()
            r1.text = f"PHASE {ph_id}"
            r1.font.name = self.fonts["mono"]
            r1.font.size = Pt(7)
            r1.font.bold = True
            r1.font.color.rgb = self.colors["brand_purple"]
            r1.font.spacing = Pt(1.5)

            # English name(衬线)
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.space_before = Pt(4)
            r2 = p2.add_run()
            r2.text = ph_name
            r2.font.name = self.fonts["serif"]
            r2.font.size = Pt(11)
            r2.font.bold = True
            r2.font.color.rgb = self.colors["text_0"]

            # Chinese name(赭石小字)
            if ph_zh:
                p3 = tf.add_paragraph()
                p3.alignment = PP_ALIGN.CENTER
                r3 = p3.add_run()
                r3.text = ph_zh
                r3.font.name = self.fonts["sans"]
                r3.font.size = Pt(8)
                r3.font.bold = True
                r3.font.color.rgb = self.colors["accent_deep"]

            # Arrow(下一个 pill 之间)
            if i < n - 1:
                arrow_x = x + (i + 1) * col_w - 0.15
                self._add_text(
                    slide, x=arrow_x - 0.05, y=y + h / 2 - 0.15, w=0.3, h=0.3,
                    text="→",
                    font=self.fonts["sans"], size=14, bold=True,
                    color=self.colors["accent"], align="center",
                )

    def _draw_delivery_box(self, slide, x, y, w, h, items):
        """交付物 box · 4 个并列卡"""
        n = len(items)
        col_w = w / n
        for i, item in enumerate(items):
            name, size = item[0], item[1] if len(item) > 1 else ""
            cx = x + i * col_w + 0.05
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cx), Inches(y), Inches(col_w - 0.1), Inches(h))
            shape.adjustments[0] = 0.1
            shape.fill.solid()
            # 金色淡背景
            shape.fill.fore_color.rgb = self.colors["bg_1"]
            sp = shape.fill.fore_color._xFill
            etree.SubElement(sp, qn("a:alpha"), {"val": "60000"})
            shape.line.color.rgb = self.colors["brand_gold"]

            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p1 = tf.paragraphs[0]
            p1.alignment = PP_ALIGN.CENTER
            r1 = p1.add_run()
            r1.text = name
            r1.font.name = self.fonts["serif"]
            r1.font.size = Pt(13)
            r1.font.bold = True
            r1.font.color.rgb = self.colors["text_0"]
            if size:
                p2 = tf.add_paragraph()
                p2.alignment = PP_ALIGN.CENTER
                r2 = p2.add_run()
                r2.text = size
                r2.font.name = self.fonts["mono"]
                r2.font.size = Pt(8)
                r2.font.color.rgb = self.colors["text_3"]

    def _draw_iron_rule_band(self, slide, x, y, w, h, text):
        """铁律横幅"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.adjustments[0] = 0.15
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors["industry_vertical"]
        sp = shape.fill.fore_color._xFill
        etree.SubElement(sp, qn("a:alpha"), {"val": "8000"})
        shape.line.color.rgb = self.colors["industry_vertical"]

        tf = shape.text_frame
        tf.margin_left = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = self.fonts["sans"]
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = self.colors["text_1"]

    def _draw_data_pain_callout(self, slide, x, y, w, h, text, color="industry_vertical"):
        """痛点 callout · 左色边"""
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(0.06), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors[color] if color in self.colors else self.colors["industry_vertical"]
        shape.line.fill.background()

        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x + 0.06), Inches(y), Inches(w - 0.06), Inches(h))
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.colors[color] if color in self.colors else self.colors["industry_vertical"]
        sp = bg.fill.fore_color._xFill
        etree.SubElement(sp, qn("a:alpha"), {"val": "6000"})
        bg.line.fill.background()

        tf = bg.text_frame
        tf.margin_left = Inches(0.2)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.name = self.fonts["sans"]
        run.font.size = Pt(10)
        run.font.color.rgb = self.colors["text_1"]

    def _tone_color_key(self, tone, fallback="accent"):
        """把语义 tone 映射到 v5 色板,避免引入新主题色。"""
        tone_map = {
            "green": "brand_green",
            "blue": "brand_blue",
            "gold": "brand_gold",
            "red": "brand_red",
            "purple": "brand_purple",
            "pink": "brand_pink",
            "accent": "accent",
            "deep": "accent_deep",
            "medical": "industry_vertical",
            "neutral": "text_2",
        }
        return tone_map.get(str(tone).lower(), fallback)

    def _draw_decision_timeline(self, slide, x, y, w, h, items):
        """时间窗 / 治疗时机 · 用 timeline 表达先后,不是 bar。"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors["bg_surface"]
        card.line.color.rgb = self.colors["border_soft"]
        self._add_card_shadow(card)

        self._add_text(
            slide, x=x + 0.3, y=y + 0.18, w=w - 0.6, h=0.3,
            text="Decision window timeline",
            font=self.fonts["mono"], size=7, bold=True,
            color=self.colors["text_3"],
        )

        n = max(len(items), 1)
        slot_w = (w - 0.8) / n
        line_y = y + 0.62
        self._add_horizontal_line(slide, x + 0.45, line_y, w - 0.9, self.colors["border"], weight=1.2)

        for i, item in enumerate(items):
            if isinstance(item, dict):
                label = item.get("label", "")
                window = item.get("window", item.get("time", ""))
                message = item.get("message", item.get("note", ""))
                tone = item.get("tone", "accent")
            else:
                label = item[0] if len(item) > 0 else ""
                window = item[1] if len(item) > 1 else ""
                message = item[2] if len(item) > 2 else ""
                tone = item[3] if len(item) > 3 else "accent"

            color_key = self._tone_color_key(tone)
            cx = x + 0.45 + i * slot_w
            marker = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(cx + slot_w * 0.5 - 0.08), Inches(line_y - 0.08),
                Inches(0.16), Inches(0.16)
            )
            marker.fill.solid()
            marker.fill.fore_color.rgb = self.colors[color_key]
            marker.line.color.rgb = self.colors["white"]
            marker.line.width = Pt(1.2)

            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(cx + 0.03), Inches(y + 0.78),
                Inches(slot_w - 0.08), Inches(h - 0.9)
            )
            box.adjustments[0] = 0.12
            box.fill.solid()
            box.fill.fore_color.rgb = self.colors[color_key]
            sp = box.fill.fore_color._xFill
            etree.SubElement(sp, qn("a:alpha"), {"val": "12000"})
            box.line.color.rgb = self.colors[color_key]
            box.line.width = Pt(0.75)

            self._add_text(
                slide, x=cx + 0.15, y=y + 0.86, w=slot_w - 0.3, h=0.18,
                text=window,
                font=self.fonts["mono"], size=7, bold=True,
                color=self.colors[color_key],
                align="center",
            )
            self._add_text(
                slide, x=cx + 0.15, y=y + 1.05, w=slot_w - 0.3, h=0.25,
                text=label,
                font=self.fonts["serif"], size=10, bold=True,
                color=self.colors["text_0"],
                align="center",
            )
            if message:
                self._add_text(
                    slide, x=cx + 0.15, y=y + 1.29, w=slot_w - 0.3, h=0.28,
                    text=message,
                    font=self.fonts["sans"], size=7,
                    color=self.colors["text_2"],
                    align="center",
                )

    def _draw_matrix(self, slide, x, y, w, h, rows):
        """机制 / 产品 / 触发条件匹配矩阵。"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors["bg_surface"]
        card.line.color.rgb = self.colors["border_soft"]
        self._add_card_shadow(card)

        self._add_text(
            slide, x=x + 0.25, y=y + 0.18, w=w - 0.5, h=0.28,
            text="Mechanism / product map",
            font=self.fonts["mono"], size=7, bold=True,
            color=self.colors["text_3"],
        )

        headers = ["Mechanism", "Clinical trigger", "Product path"]
        col_ws = [w * 0.30, w * 0.33, w * 0.37]
        header_y = y + 0.55
        row_start = y + 0.95
        row_h = (h - 1.15) / max(len(rows), 1)
        cur_x = x + 0.25
        for ci, header in enumerate(headers):
            self._add_text(
                slide, x=cur_x, y=header_y, w=col_ws[ci] - 0.12, h=0.25,
                text=header,
                font=self.fonts["sans"], size=7, bold=True,
                color=self.colors["accent_deep"],
            )
            cur_x += col_ws[ci]

        for ri, row in enumerate(rows):
            if isinstance(row, dict):
                mechanism = row.get("mechanism", "")
                trigger = row.get("trigger", "")
                path = row.get("path", row.get("product", ""))
                tone = row.get("tone", "accent")
            else:
                mechanism = row[0] if len(row) > 0 else ""
                trigger = row[1] if len(row) > 1 else ""
                path = row[2] if len(row) > 2 else ""
                tone = row[3] if len(row) > 3 else "accent"

            color_key = self._tone_color_key(tone)
            cy = row_start + ri * row_h
            cur_x = x + 0.25
            values = [mechanism, trigger, path]
            for ci, value in enumerate(values):
                cell = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    Inches(cur_x), Inches(cy),
                    Inches(col_ws[ci] - 0.12), Inches(row_h - 0.08)
                )
                cell.adjustments[0] = 0.08
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.colors["bg_1"]
                cell.line.color.rgb = self.colors["border_soft"]

                if ci == 0:
                    band = slide.shapes.add_shape(
                        MSO_SHAPE.RECTANGLE,
                        Inches(cur_x), Inches(cy),
                        Inches(0.06), Inches(row_h - 0.08)
                    )
                    band.fill.solid()
                    band.fill.fore_color.rgb = self.colors[color_key]
                    band.line.fill.background()

                self._add_text(
                    slide, x=cur_x + 0.12, y=cy + 0.08,
                    w=col_ws[ci] - 0.36, h=row_h - 0.16,
                    text=str(value),
                    font=self.fonts["sans"], size=8,
                    color=self.colors["text_1"],
                    line_spacing=1.15,
                )
                cur_x += col_ws[ci]

    def _draw_evidence_list(self, slide, x, y, w, h, items):
        """证据链列表 · 用分层文本/标记表达,不是通过率条。"""
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h)
        )
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = self.colors["bg_surface"]
        card.line.color.rgb = self.colors["border_soft"]
        self._add_card_shadow(card)

        self._add_text(
            slide, x=x + 0.25, y=y + 0.18, w=w - 0.5, h=0.28,
            text="Evidence chain",
            font=self.fonts["mono"], size=7, bold=True,
            color=self.colors["text_3"],
        )

        row_h = (h - 0.75) / max(len(items), 1)
        for i, item in enumerate(items):
            if isinstance(item, dict):
                label = item.get("label", "")
                note = item.get("note", "")
                grade = item.get("grade", item.get("level", "A"))
                tone = item.get("tone", "green")
            else:
                label = item[0] if len(item) > 0 else ""
                note = item[1] if len(item) > 1 else ""
                grade = item[2] if len(item) > 2 else "A"
                tone = item[3] if len(item) > 3 else "green"

            color_key = self._tone_color_key(tone)
            cy = y + 0.6 + i * row_h
            chip = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + 0.25), Inches(cy + 0.06),
                Inches(0.42), Inches(0.28)
            )
            chip.adjustments[0] = 0.4
            chip.fill.solid()
            chip.fill.fore_color.rgb = self.colors[color_key]
            chip.line.fill.background()
            tf = chip.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = str(grade)
            run.font.name = self.fonts["mono"]
            run.font.size = Pt(7)
            run.font.bold = True
            run.font.color.rgb = self.colors["white"]

            self._add_text(
                slide, x=x + 0.82, y=cy, w=w - 1.1, h=0.25,
                text=str(label),
                font=self.fonts["serif"], size=10, bold=True,
                color=self.colors["text_0"],
            )
            if note:
                self._add_text(
                    slide, x=x + 0.82, y=cy + 0.25, w=w - 1.1, h=row_h - 0.24,
                    text=str(note),
                    font=self.fonts["sans"], size=7,
                    color=self.colors["text_2"],
                    line_spacing=1.15,
                )

    def _draw_proof_bar(self, slide, x, y, w, h, name, domain, pct, status):
        """横向校验条"""
        # 名称(180px)
        self._add_text(
            slide, x=x, y=y, w=2.0, h=h,
            text=name,
            font=self.fonts["serif"], size=11, bold=True,
            color=self.colors["text_0"],
        )
        if domain:
            self._add_text(
                slide, x=x, y=y + 0.25, w=2.0, h=h,
                text=domain,
                font=self.fonts["sans"], size=7,
                color=self.colors["text_3"],
            )
        # 进度条 track
        track_x = x + 2.2
        track_w = w - 3.0
        track = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(track_x), Inches(y + 0.1), Inches(track_w), Inches(h - 0.2))
        track.adjustments[0] = 0.5
        track.fill.solid()
        track.fill.fore_color.rgb = self.colors["bg_1"]
        track.line.fill.background()
        # fill
        try:
            pct_val = float(pct.rstrip('%'))
        except (ValueError, AttributeError):
            pct_val = 50
        fill_w = (track_w * pct_val / 100)
        is_pass = pct_val >= 75
        fill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(track_x), Inches(y + 0.1), Inches(fill_w), Inches(h - 0.2))
        fill.adjustments[0] = 0.5
        fill.fill.solid()
        fill.fill.fore_color.rgb = self.colors["brand_green"] if is_pass else self.colors["brand_gold"]
        fill.line.fill.background()
        tf = fill.text_frame
        tf.margin_left = Inches(0.15)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = pct
        run.font.name = self.fonts["serif"]
        run.font.size = Pt(10)
        run.font.bold = True
        run.font.color.rgb = self.colors["white"]
        # status
        self._add_text(
            slide, x=x + w - 0.8, y=y, w=0.8, h=h,
            text=status,
            font=self.fonts["sans"], size=8, bold=True,
            color=self.colors["brand_green"] if is_pass else self.colors["brand_gold"],
            align="right",
        )

    def _draw_tech_stack_row(self, slide, x, y, w, h, items):
        """tech stack 胶囊行"""
        n = len(items)
        item_w = w / n
        for i, name in enumerate(items):
            shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                Inches(x + i * item_w + 0.05), Inches(y), Inches(item_w - 0.1), Inches(h))
            shape.adjustments[0] = 0.5
            shape.fill.solid()
            shape.fill.fore_color.rgb = self.colors["bg_surface"]
            shape.line.color.rgb = self.colors["border_soft"]
            tf = shape.text_frame
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = name
            run.font.name = self.fonts["mono"]
            run.font.size = Pt(8)
            run.font.bold = True
            run.font.color.rgb = self.colors["text_2"]

    def _draw_cta_card(self, slide, x, y, w, h, tag, time, title, desc, primary=True):
        """CTA card · 主钩 / 辅钩"""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y), Inches(w), Inches(h))
        shape.adjustments[0] = 0.06
        shape.fill.solid()
        shape.fill.fore_color.rgb = self.colors["bg_surface"]
        if primary:
            shape.line.color.rgb = self.colors["accent"]
            shape.line.width = Pt(1.5)
        else:
            shape.line.color.rgb = self.colors["border_soft"]
            shape.line.dash_style = 7  # dash
        self._add_card_shadow(shape)

        tf = shape.text_frame
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.2)

        # tag
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = tag
        r1.font.name = self.fonts["sans"]
        r1.font.size = Pt(7)
        r1.font.bold = True
        r1.font.color.rgb = self.colors["accent"] if primary else self.colors["text_2"]
        r1.font.spacing = Pt(1.2)

        # time
        p2 = tf.add_paragraph()
        p2.space_before = Pt(8)
        r2 = p2.add_run()
        r2.text = time
        r2.font.name = self.fonts["serif"]
        r2.font.size = Pt(20)
        r2.font.bold = True
        r2.font.color.rgb = self.colors["accent"] if primary else self.colors["text_2"]

        # title
        p3 = tf.add_paragraph()
        p3.space_before = Pt(8)
        r3 = p3.add_run()
        r3.text = title
        r3.font.name = self.fonts["serif"]
        r3.font.size = Pt(13)
        r3.font.bold = True
        r3.font.color.rgb = self.colors["text_0"]

        # desc
        p4 = tf.add_paragraph()
        p4.space_before = Pt(6)
        p4.line_spacing = 1.4
        r4 = p4.add_run()
        r4.text = desc
        r4.font.name = self.fonts["sans"]
        r4.font.size = Pt(9)
        r4.font.color.rgb = self.colors["text_2"]

    def _add_jump_row(self, slide, x, y, links):
        """跳转链接行"""
        for i, (label, url) in enumerate(links):
            self._add_text(
                slide, x=x + i * 2.3, y=y, w=2.2, h=0.3,
                text=f"↗ {label}",
                font=self.fonts["mono"], size=9, bold=True,
                color=self.colors["accent_deep"],
            )

    def _add_horizontal_line(self, slide, x, y, w, color, weight=0.75):
        """水平细线"""
        line = slide.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
        line.line.color.rgb = color
        line.line.width = Pt(weight)

    def _add_card_shadow(self, shape):
        """给 shape 添加软阴影"""
        sp = shape.shadow
        sp.inherit = False
        # 使用 OOXML 直接添加
        sp_pr = shape._element.spPr
        effectLst = etree.SubElement(sp_pr, qn("a:effectLst"))
        outerShdw = etree.SubElement(effectLst, qn("a:outerShdw"), {
            "blurRad": "180000",  # blur 半径
            "dist": "60000",       # 距离
            "dir": "5400000",      # 方向(90 度,向下)
            "rotWithShape": "0",
        })
        clr = etree.SubElement(outerShdw, qn("a:srgbClr"), {"val": "141413"})
        alpha = etree.SubElement(clr, qn("a:alpha"), {"val": "8000"})


# ============================================================
# 命令行测试入口
# ============================================================

if __name__ == "__main__":
    """演示 · 生成血液科 v5 案例的 PPT 镜像"""
    deck = EditorialDeck(
        title="血液科 IFI · v5",
        industry="medical",
        embed_fonts=False,  # 测试用降级字体
    )

    # Slide 1: HERO
    deck.add_hero_slide(
        eyebrow="HEMATOLOGY · IFI MARKET INTELLIGENCE · 2026",
        title_main="一句话起跑 ",
        title_accent="AI 编排血液真菌市场调研",
        sub="用户输入一句话疾病名 → 系统跑通 5 phase 流水线 → 输出 4 件交付物",
        stats=[
            ("38亿", "总市场", "7×6×双路径"),
            ("7", "Modules", "独立子页"),
            ("22", "Charts", "1+6+15"),
            ("43", "PMID", "可点击"),
        ],
    )

    # Slide 2: WHY
    deck.add_why_slide(
        eyebrow="WHY · 项目源头",
        title="血液真菌市场 = ",
        title_accent="两条平行患者旅程",
        bar_chart={
            "axis": ["", "预防", "突破", "经验性", "诊断驱动", "目标治疗", "维持"],
            "rows": [
                {
                    "label": "LINE A · 有预防", "pct": "~50%",
                    "cells": [
                        {"type": "s-prevent", "text": "预防"},
                        {"type": "event", "text": "⚡突破"},
                        {"type": "s-empiric", "text": "经验性"},
                        {"type": "s-dd", "text": "诊断驱动"},
                        {"type": "s-target", "text": "目标治疗"},
                        {"type": "s-maint", "text": "维持"},
                    ],
                },
                {
                    "label": "LINE B · 无预防", "pct": "~50%",
                    "cells": [
                        {"type": "empty"}, {"type": "empty"},
                        {"type": "s-empiric", "text": "经验性"},
                        {"type": "s-dd", "text": "诊断驱动"},
                        {"type": "s-target", "text": "目标治疗"},
                        {"type": "s-maint", "text": "维持"},
                    ],
                },
            ],
        },
        data_pain="B 从 1% → 15% 不等(差 12 倍)· 销售收集准确率不足 · 必须用 AI 算清 P × B 真实分布",
    )

    # Slide 3: EVIDENCE
    deck.add_evidence_slide(
        eyebrow="CITE-OR-BLOCK · 语义层",
        title="每一个数字 · 都能点击回溯到原文",
        funnel=[
            ("PubMed 触达", "数千"),
            ("初筛", "~500"),
            ("全文+时间窗", "~120"),
            ("入库", "47"),
            ("纳入 A/B/C", "43"),
        ],
        pyramid=[
            ("🟢 最高级", 18, "green"),
            ("🔵 重要级", 17, "blue"),
            ("🟡 参考级", 7, "gold"),
            ("⚪ 国际", 2, "neutral"),
            ("🟠 估算", "按需", "deep"),
        ],
    )

    # Slide 4: RELATIONSHIPS
    deck.add_relationship_slide(
        eyebrow="RELATIONSHIP GRAMMAR · 非条形图默认",
        title="时间窗、机制匹配、证据链 · 分开表达",
        note="独立 KPI 用 stat;治疗时机用 timeline;机制/产品匹配用 matrix;同一指标比较才用 proof bars。",
        timeline=[
            {"label": "Risk signal", "window": "0-6h", "message": "先识别治疗窗口", "tone": "gold"},
            {"label": "Decision", "window": "6-24h", "message": "明确是否升级", "tone": "accent"},
            {"label": "Target", "window": "24-48h", "message": "按证据校准路径", "tone": "green"},
            {"label": "Review", "window": "48h+", "message": "回看适应证和证据", "tone": "blue"},
        ],
        matrix=[
            {"mechanism": "耐药风险", "trigger": "高危病原线索 + 既往用药", "path": "产品路径 A:早期覆盖假设", "tone": "medical"},
            {"mechanism": "治疗时机", "trigger": "ICU / 免疫抑制 / 失败信号", "path": "产品路径 B:窗口前移", "tone": "accent"},
            {"mechanism": "证据锚点", "trigger": "指南 / RCT / 本地真实世界", "path": "产品路径 C:可解释落地", "tone": "green"},
        ],
        evidence=[
            {"grade": "A", "label": "指南锚点", "note": "用于定义必须遵守的临床边界", "tone": "green"},
            {"grade": "B", "label": "研究证据", "note": "用于支持机制和治疗路径", "tone": "blue"},
            {"grade": "C", "label": "市场访谈", "note": "用于解释行为阻力和落地动作", "tone": "gold"},
        ],
    )

    # Slide 5: ARCHITECTURE
    deck.add_architecture_slide(
        eyebrow="THE ARCHITECTURE",
        title="一句话进 → 4 件交付物出",
        phases=[
            ("0", "Contract", "契约抽取"),
            ("1", "Stratify", "动态分层"),
            ("2", "Recall", "数据召回"),
            ("3", "Compose", "内容生成"),
            ("4", "Verify", "事实核验"),
            ("5", "Bundle", "打包交付"),
        ],
        deliverables=[
            ("📄 HTML", "9MB"),
            ("📑 PDF", "3 份"),
            ("📊 Excel", "7 张"),
            ("🗂 Manifest", ".json"),
        ],
        iron_rule="P0 IRON RULE · Cite-or-Block · 禁穷举 / 禁凭记忆 / 禁丢原文 · 每事实声明带锚点",
    )

    # Slide 6: PROOF
    deck.add_proof_slide(
        eyebrow="PROOF · SAME METRIC ONLY",
        title="同一指标的横向通过率比较",
        proof_bars=[
            {"name": "ALK NSCLC v2", "domain": "肺癌靶向", "pct": "96.7%", "status": "PASS"},
            {"name": "原发性高血压", "domain": "心血管慢病", "pct": "95.0%", "status": "PASS"},
            {"name": "血友病",       "domain": "罕见血液病",   "pct": "77.3%", "status": "PASS"},
            {"name": "GLP-1 RA T2DM", "domain": "内分泌慢病", "pct": "67.5%", "status": "VERIFY"},
            {"name": "HER2 乳腺癌",  "domain": "肿瘤抗体偶联", "pct": "21.6%", "status": "VERIFY"},
        ],
        key_message="关键事实:5 个疾病不管通过率多少,4 件交付物全跑通 (HTML / PDF / XLSX / Manifest)",
    )

    # Slide 7: CTA
    deck.add_cta_slide(
        eyebrow="THE TEAM & CALL-TO-ACTION",
        team=["YongQi", "SimonSu", "VivienZhan", "RuiYu", "YingJi"],
        tech_stack=["Claude Code", "22 Skills", "Codex GPT-5.5", "Mermaid ELK", "Puppeteer", "Cite-or-Block"],
        ctas=[
            {"tag": "主钩 A · 轻量门票", "time": "1 周", "title": "一句话疾病名 → demo",
             "desc": "给我一句话,1 周内 demo 给你看。基于 Phase 2a 已验证的 5 phase 流水线。", "primary": True},
            {"tag": "主钩 B · 深度合作", "time": "1 月", "title": "BU GPT 助手 · 架构底座",
             "desc": "BU 内部 GPT 助手基于这套架构(22 Skill + Cite-or-Block + 3 道质检门)做底座。", "primary": True},
            {"tag": "Q&A 加码", "time": "3 月", "title": "5 治疗领域同步",
             "desc": "肿瘤 / 心血管 / 内分泌 / 血液 / 罕见病 — 3 月 5 套报告同步出。", "primary": False},
        ],
        contact="yong.qi.gpt@gmail.com",
    )

    out = deck.save("editorial_demo.pptx")
    print(f"Generated: {out}")
