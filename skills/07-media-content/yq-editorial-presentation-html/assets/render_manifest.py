#!/usr/bin/env python3
"""Render one validated editorial manifest to HTML and/or PPTX."""

from __future__ import annotations

import argparse
import hashlib
import html
import importlib.util
import json
from pathlib import Path
from typing import Any

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
SHELL_PATH = Path(__file__).with_name("deck-shell.html")
PLANNER_PATH = Path(__file__).with_name("plan_deck.py")
GENERATOR_PATH = Path(__file__).with_name("generate_pptx.py")


def _load_module(name: str, path: Path):
    spec = importlib.util.spec_from_file_location(name, path)
    if spec is None or spec.loader is None:
        raise RuntimeError(f"Unable to load module: {path}")
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def _normalized_manifest(manifest: dict[str, Any]) -> bytes:
    return json.dumps(manifest, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode("utf-8")


def _manifest_hash(manifest: dict[str, Any]) -> str:
    return hashlib.sha256(_normalized_manifest(manifest)).hexdigest()


def _items(slide: dict[str, Any], count: int = 4) -> list[str]:
    values = slide.get("items") or slide.get("notes") or []
    normalized = [str(item.get("label", item)) if isinstance(item, dict) else str(item) for item in values]
    return (normalized + [f"Evidence {index + 1}" for index in range(count)])[:count]


def _html_body(slide: dict[str, Any]) -> tuple[str, list[str]]:
    comp = slide["composition_id"]
    title = html.escape(str(slide.get("title", slide["id"])))
    body = html.escape(str(slide.get("body", "Evidence-led editorial narrative.")))
    metric = html.escape(str(slide.get("metric", slide.get("value", "64%"))))
    items = [html.escape(item) for item in _items(slide, 6)]

    if comp == "statement-full-bleed":
        return f'<div class="statement"><h1>{title}</h1><p>{body}</p></div>', ["statement"]
    if comp == "dominant-metric-marginalia":
        notes = "".join(f"<li>{item}</li>" for item in items[:3])
        return f'<div class="metric-anchor"><strong>{metric}</strong><span>{title}</span></div><ul class="marginalia">{notes}</ul>', ["metric-anchor", "marginalia"]
    if comp == "chart-notes-asymmetric":
        bars = "".join(f'<i style="--v:{35 + index * 12}%"></i>' for index in range(4))
        notes = "".join(f"<li>{item}</li>" for item in items[:4])
        return f'<div class="chart-field"><h2>{title}</h2><div class="bars">{bars}</div></div><ol class="evidence-notes">{notes}</ol>', ["chart-field", "evidence-notes"]
    if comp == "comparison-split":
        return f'<article class="compare-a"><small>BEFORE</small><h2>{items[0]}</h2><p>{body}</p></article><article class="compare-b"><small>AFTER</small><h2>{items[1]}</h2><p>{body}</p></article>', ["compare-a", "compare-b"]
    if comp == "ledger-takeaway-band":
        rows = "".join(f"<div><span>{index + 1:02d}</span><b>{item}</b><em>VERIFIED</em></div>" for index, item in enumerate(items[:5]))
        return f'<div class="ledger"><h2>{title}</h2>{rows}</div><aside class="takeaway">{body}</aside>', ["ledger", "takeaway"]
    if comp == "process-rail":
        stages = "".join(f"<div><b>{index + 1:02d}</b><span>{item}</span></div>" for index, item in enumerate(items[:5]))
        return f'<h2 class="rail-title">{title}</h2><div class="process-rail">{stages}</div>', ["rail-title", "process-rail"]
    if comp == "matrix-marginalia":
        cells = "".join(f"<span>{items[index % len(items)]}</span>" for index in range(9))
        return f'<div class="matrix"><h2>{title}</h2><div>{cells}</div></div><aside class="matrix-note">{body}</aside>', ["matrix", "matrix-note"]
    if comp == "small-multiple-field":
        panels = "".join(f'<article><b>{item}</b><i style="--v:{40 + index * 8}%"></i></article>' for index, item in enumerate(items[:6]))
        return f'<h2 class="multiple-title">{title}</h2><div class="small-multiples">{panels}</div>', ["multiple-title", "small-multiples"]
    if comp == "image-text-offset":
        return f'<div class="media-anchor" role="img" aria-label="{title}"><span>{metric}</span></div><article class="media-copy"><h2>{title}</h2><p>{body}</p></article>', ["media-anchor", "media-copy"]
    tiles = "".join(f'<article class="tile-{index + 1}"><b>{item}</b><span>{metric if index == 0 else body}</span></article>' for index, item in enumerate(items[:5]))
    return f'<h2 class="mosaic-title">{title}</h2><div class="mosaic">{tiles}</div>', ["mosaic-title", "mosaic"]


COMPOSITION_CSS = r"""
  .manifest-slide{background:var(--canvas);color:var(--ink)}
  .manifest-slide[data-tone-id="surface"]{background:var(--surface)}
  .manifest-slide[data-tone-id="contrast_surface"]{background:color-mix(in srgb,var(--canvas) 82%,var(--primary))}
  .manifest-slide[data-tone-id="statement"]{background:color-mix(in srgb,var(--canvas) 70%,var(--secondary))}
  .manifest-slide .content{display:grid;grid-template-columns:repeat(12,1fr);grid-template-rows:repeat(8,1fr);gap:clamp(10px,1vw,22px);height:100%;min-height:0;margin-top:clamp(24px,3vh,44px)}
  .manifest-slide h1,.manifest-slide h2,.manifest-slide p{margin:0}.manifest-slide h1,.manifest-slide h2{font-family:var(--font-display,Georgia,serif)}
  .statement{grid-column:1/13;grid-row:1/9;display:flex;flex-direction:column;justify-content:center;max-width:1200px}.statement h1{font-size:clamp(86px,7.5vw,144px);line-height:1}.statement p{max-width:800px;margin-top:30px;font-size:clamp(25px,2vw,38px);color:var(--muted)}
  .metric-anchor{grid-column:1/8;grid-row:1/9;display:flex;flex-direction:column;justify-content:center}.metric-anchor strong{font:700 clamp(130px,14vw,260px)/.8 var(--font-display,Georgia,serif);color:var(--primary)}.metric-anchor span{margin-top:32px;font-size:clamp(30px,3vw,56px)}.marginalia{grid-column:9/13;grid-row:2/8;margin:0;padding:0;display:grid;align-content:center;gap:22px;list-style:none}.marginalia li{padding:20px 0;border-top:2px solid var(--secondary);font-size:clamp(18px,1.4vw,28px)}
  .chart-field{grid-column:1/9;grid-row:1/9;padding:32px;background:var(--surface);border-left:8px solid var(--primary)}.chart-field h2{font-size:clamp(52px,4.5vw,88px)}.bars{height:65%;display:flex;align-items:end;gap:8%;margin-top:7%}.bars i{flex:1;height:var(--v);background:var(--secondary)}.evidence-notes{grid-column:10/13;grid-row:2/8;margin:0;padding-left:1.4em;display:grid;align-content:center;gap:20px;font-size:clamp(17px,1.25vw,25px)}
  .compare-a,.compare-b{grid-row:1/9;padding:clamp(28px,3vw,54px);display:flex;flex-direction:column;justify-content:center}.compare-a{grid-column:1/6;background:var(--surface);border:2px solid color-mix(in srgb,var(--ink) 18%,transparent)}.compare-b{grid-column:6/13;background:var(--contrast-surface);color:var(--canvas)}.compare-a h2,.compare-b h2{margin:20px 0;font-size:clamp(48px,4vw,78px)}.compare-a p,.compare-b p{font-size:clamp(18px,1.4vw,27px)}
  .ledger{grid-column:1/13;grid-row:1/7}.ledger h2{font-size:clamp(48px,4vw,78px);margin-bottom:20px}.ledger>div{display:grid;grid-template-columns:80px 1fr 160px;gap:20px;padding:13px 8px;border-top:1px solid color-mix(in srgb,var(--ink) 20%,transparent);font-size:clamp(16px,1.2vw,23px)}.ledger em{font-family:var(--font-data,monospace);color:var(--positive)}.takeaway{grid-column:1/13;grid-row:7/9;padding:24px 32px;background:var(--primary);color:var(--surface);font-size:clamp(20px,1.6vw,30px)}
  .rail-title{grid-column:1/13;grid-row:1/3;font-size:clamp(58px,5vw,92px)}.process-rail{grid-column:1/13;grid-row:4/8;display:grid;grid-template-columns:repeat(5,1fr);position:relative}.process-rail::before{content:"";position:absolute;left:8%;right:8%;top:34%;height:4px;background:var(--primary)}.process-rail>div{z-index:1;display:flex;flex-direction:column;align-items:center;text-align:center;gap:22px}.process-rail b{display:grid;place-items:center;width:74px;height:74px;border-radius:50%;background:var(--canvas);border:5px solid var(--primary);font-family:var(--font-data,monospace)}.process-rail span{font-size:clamp(16px,1.2vw,23px)}
  .matrix{grid-column:1/9;grid-row:1/9}.matrix h2{font-size:clamp(50px,4.2vw,82px);margin-bottom:18px}.matrix>div{height:72%;display:grid;grid-template-columns:repeat(3,1fr);gap:8px}.matrix span{display:grid;place-items:center;padding:12px;background:var(--surface);border:1px solid color-mix(in srgb,var(--ink) 18%,transparent);text-align:center}.matrix-note{grid-column:10/13;grid-row:2/8;padding:28px;border-left:6px solid var(--secondary);font-size:clamp(18px,1.35vw,26px)}
  .multiple-title{grid-column:1/13;grid-row:1/3;font-size:clamp(52px,4.5vw,86px)}.small-multiples{grid-column:1/13;grid-row:3/9;display:grid;grid-template-columns:repeat(3,1fr);grid-template-rows:repeat(2,1fr);gap:14px}.small-multiples article{position:relative;padding:20px;background:var(--surface);border-top:5px solid var(--secondary);overflow:hidden}.small-multiples i{position:absolute;left:20px;right:20px;bottom:20px;height:var(--v);background:color-mix(in srgb,var(--primary) 45%,transparent)}
  .media-anchor{grid-column:1/8;grid-row:1/9;display:grid;place-items:center;background:var(--primary);color:var(--surface)}.media-anchor span{font:700 clamp(100px,10vw,190px)/1 var(--font-display,Georgia,serif)}.media-copy{grid-column:9/13;grid-row:2/8;display:flex;flex-direction:column;justify-content:center}.media-copy h2{font-size:clamp(48px,4vw,78px)}.media-copy p{margin-top:24px;font-size:clamp(18px,1.4vw,27px);color:var(--muted)}
  .mosaic-title{grid-column:1/5;grid-row:1/3;font-size:clamp(50px,4.2vw,82px)}.mosaic{grid-column:1/13;grid-row:3/9;display:grid;grid-template-columns:2fr 1fr 1fr;grid-template-rows:1fr 1fr;gap:12px}.mosaic article{padding:22px;display:flex;flex-direction:column;justify-content:space-between;background:var(--surface);border:1px solid color-mix(in srgb,var(--ink) 16%,transparent)}.mosaic .tile-1{grid-row:1/3;background:var(--contrast-surface);color:var(--canvas)}.mosaic .tile-4{grid-column:2/4}.mosaic article b{font-size:clamp(18px,1.4vw,28px)}
"""


def render_html(manifest: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    profile = manifest["palette"]
    profile_css = ":root{" + ";".join(
        f"--{key.replace('_', '-')}:{value}" for key, value in profile.items()
    ) + "}"
    sections = []
    signatures = []
    total = len(manifest["slides"])
    for index, slide in enumerate(manifest["slides"], 1):
        body, regions = _html_body(slide)
        sections.append(
            f'<section class="slide manifest-slide composition-{html.escape(slide["composition_id"])}" '
            f'data-slide-id="{html.escape(slide["id"])}" data-composition-id="{html.escape(slide["composition_id"])}" data-tone-id="{html.escape(slide["tone_id"])}">'
            f'<div class="slide-chrome"><span>{html.escape(str(slide.get("eyebrow", slide.get("narrative_role", "CONTENT"))).upper())}</span><span>{index:02d} / {total:02d}</span></div>'
            f'<div class="content">{body}</div>'
            f'<footer class="slide-foot"><span>{html.escape(slide["expression_family"])}</span><span>{html.escape(slide["tone_id"])}</span></footer></section>'
        )
        signatures.append({"id": slide["id"], "composition_id": slide["composition_id"], "regions": regions})

    shell = SHELL_PATH.read_text(encoding="utf-8")
    marker = "    <!-- Insert generated <section class=\"slide\"> elements from the validated manifest here. -->"
    if marker not in shell:
        raise RuntimeError("Deck shell insertion marker missing")
    shell = shell.replace(marker, "\n".join(sections))
    digest = _manifest_hash(manifest)
    shell = shell.replace("</head>", f'<meta name="manifest-sha256" content="{digest}"><style id="profile-and-compositions">{profile_css}{COMPOSITION_CSS}</style></head>')
    output_path.parent.mkdir(parents=True, exist_ok=True)
    output_path.write_text(shell, encoding="utf-8")
    return signatures


def _pptx_color(deck, key: str):
    aliases = {"canvas": "canvas", "surface": "surface", "contrast": "contrast_surface", "primary": "primary", "secondary": "secondary"}
    return deck.colors[aliases.get(key, key)]


def _pptx_box(slide, deck, bounds, fill="surface", line="border_soft", radius=False):
    x, y, w, h = bounds
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = _pptx_color(deck, fill)
    if line:
        shape.line.color.rgb = deck.colors.get(line, _pptx_color(deck, "ink")); shape.line.width = Pt(0.8)
    else:
        shape.line.fill.background()
    return shape


def _pptx_text(slide, deck, bounds, text, size=20, font="sans", color="text_0", bold=False, align="left"):
    x, y, w, h = bounds
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame; frame.clear(); frame.word_wrap = True; frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]; paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[align]
    run = paragraph.add_run(); run.text = str(text); run.font.name = deck.fonts[font]; run.font.size = Pt(size); run.font.bold = bold; run.font.color.rgb = deck.colors.get(color, deck.colors["text_0"])
    return box


def _pptx_regions(slide_spec: dict[str, Any]) -> list[tuple[float, float, float, float]]:
    comp = slide_spec["composition_id"]
    return {
        "statement-full-bleed": [(0.8, 1.5, 11.7, 4.8)],
        "dominant-metric-marginalia": [(0.7, 1.5, 7.2, 4.9), (8.6, 1.7, 4.0, 4.5)],
        "chart-notes-asymmetric": [(0.7, 1.6, 8.2, 5.0), (9.4, 1.8, 3.2, 4.6)],
        "comparison-split": [(0.7, 1.5, 5.0, 5.2), (5.9, 1.5, 6.7, 5.2)],
        "ledger-takeaway-band": [(0.7, 1.4, 11.9, 4.2), (0.7, 5.8, 11.9, 0.9)],
        "process-rail": [(0.7, 1.3, 11.9, 1.2), (0.7, 3.0, 11.9, 3.0)],
        "matrix-marginalia": [(0.7, 1.4, 8.3, 5.2), (9.4, 1.8, 3.2, 4.5)],
        "small-multiple-field": [(0.7, 1.3, 11.9, 1.0), (0.7, 2.5, 11.9, 4.3)],
        "image-text-offset": [(0.7, 1.3, 7.5, 5.5), (8.7, 1.8, 3.9, 4.6)],
        "editorial-mosaic": [(0.7, 1.2, 4.2, 1.1), (0.7, 2.5, 11.9, 4.2)],
    }[comp]


def _render_pptx_slide(deck, spec: dict[str, Any]):
    slide = deck.prs.slides.add_slide(deck._blank_layout)
    tone_background = {"canvas": "bg_0", "surface": "bg_surface", "contrast_surface": "bg_1", "statement": "bg_2"}
    deck._fill_background(slide, tone_background.get(spec.get("tone_id"), "bg_0"))
    comp = spec["composition_id"]
    title = str(spec.get("title", spec["id"]))
    body = str(spec.get("body", "Evidence-led editorial narrative."))
    metric = str(spec.get("metric", spec.get("value", "64%")))
    items = _items(spec, 6)
    regions = _pptx_regions(spec)
    deck._add_eyebrow_tag(slide, 0.7, 0.55, str(spec.get("eyebrow", spec.get("narrative_role", "CONTENT"))).upper())

    if comp == "statement-full-bleed":
        _pptx_text(slide, deck, regions[0], title, 54, "serif", "text_0", True)
        _pptx_text(slide, deck, (0.9, 5.7, 8.8, 0.7), body, 20, "sans", "text_2")
    elif comp == "dominant-metric-marginalia":
        _pptx_text(slide, deck, regions[0], metric, 96, "serif", "accent", True)
        _pptx_text(slide, deck, (0.9, 5.8, 7.0, 0.7), title, 25, "sans", "text_0", True)
        for index, item in enumerate(items[:3]): _pptx_text(slide, deck, (8.7, 1.8 + index * 1.35, 3.7, 1.0), item, 18, "sans", "text_0", index == 0)
    elif comp == "chart-notes-asymmetric":
        _pptx_box(slide, deck, regions[0], "surface", "border_soft")
        _pptx_text(slide, deck, (0.95, 1.7, 7.5, 0.8), title, 30, "serif", "text_0", True)
        for index in range(4): _pptx_box(slide, deck, (1.1 + index * 1.7, 5.8 - index * 0.45, 0.85, 0.7 + index * 0.45), "secondary", None)
        for index, item in enumerate(items[:4]): _pptx_text(slide, deck, (9.5, 1.8 + index * 1.05, 3.0, 0.8), item, 15, "sans", "text_0")
    elif comp == "comparison-split":
        _pptx_box(slide, deck, regions[0], "surface", "border_soft")
        _pptx_box(slide, deck, regions[1], "primary", None)
        _pptx_text(slide, deck, (1.0, 2.0, 4.4, 1.2), items[0], 28, "serif", "text_0", True)
        _pptx_text(slide, deck, (6.3, 2.0, 5.8, 1.2), items[1], 32, "serif", "white", True)
    elif comp == "ledger-takeaway-band":
        _pptx_text(slide, deck, (0.8, 1.0, 11.5, 0.7), title, 30, "serif", "text_0", True)
        for index, item in enumerate(items[:5]):
            _pptx_box(slide, deck, (0.8, 1.8 + index * 0.72, 11.5, 0.62), "surface", "border_soft")
            _pptx_text(slide, deck, (1.0, 1.8 + index * 0.72, 10.9, 0.62), f"{index + 1:02d}   {item}", 15, "mono", "text_0")
        _pptx_box(slide, deck, regions[1], "primary", None); _pptx_text(slide, deck, regions[1], body, 17, "sans", "white", True, "center")
    elif comp == "process-rail":
        _pptx_text(slide, deck, regions[0], title, 34, "serif", "text_0", True)
        _pptx_box(slide, deck, (1.0, 4.0, 11.0, 0.05), "primary", None)
        for index, item in enumerate(items[:5]):
            x = 1.0 + index * 2.6; _pptx_box(slide, deck, (x, 3.55, 0.75, 0.75), "canvas", "accent", True); _pptx_text(slide, deck, (x - 0.35, 4.5, 1.45, 1.0), item, 13, "sans", "text_0", False, "center")
    elif comp == "matrix-marginalia":
        _pptx_text(slide, deck, (0.8, 1.0, 8.0, 0.7), title, 30, "serif", "text_0", True)
        for row in range(3):
            for col in range(3):
                bounds=(0.8 + col * 2.65, 1.9 + row * 1.45, 2.45, 1.2); _pptx_box(slide, deck, bounds, "surface", "border_soft"); _pptx_text(slide, deck, bounds, items[(row * 3 + col) % len(items)], 13, "sans", "text_0", False, "center")
        _pptx_text(slide, deck, regions[1], body, 18, "sans", "text_0")
    elif comp == "small-multiple-field":
        _pptx_text(slide, deck, regions[0], title, 32, "serif", "text_0", True)
        for index, item in enumerate(items[:6]):
            col=index%3; row=index//3; bounds=(0.8+col*4.0,2.2+row*2.2,3.7,1.9); _pptx_box(slide, deck,bounds,"surface","border_soft"); _pptx_text(slide, deck,(bounds[0]+0.2,bounds[1]+0.1,3.2,0.5),item,14,"sans","text_0",True); _pptx_box(slide,deck,(bounds[0]+0.2,bounds[1]+1.1,2.5+index*.12,0.28),"secondary",None)
    elif comp == "image-text-offset":
        _pptx_box(slide, deck, regions[0], "primary", None); _pptx_text(slide, deck, regions[0], metric, 76, "serif", "white", True, "center"); _pptx_text(slide, deck, (8.8, 2.0, 3.5, 1.4), title, 32, "serif", "text_0", True); _pptx_text(slide, deck, (8.8, 3.6, 3.5, 1.8), body, 17, "sans", "text_2")
    else:
        _pptx_text(slide, deck, regions[0], title, 30, "serif", "text_0", True)
        tiles=[(0.8,2.4,5.3,4.0),(6.3,2.4,2.9,1.8),(9.4,2.4,2.9,1.8),(6.3,4.4,6.0,2.0)]
        for index,bounds in enumerate(tiles): _pptx_box(slide,deck,bounds,"primary" if index==0 else "surface",None if index==0 else "border_soft"); _pptx_text(slide,deck,bounds,items[index],19,"sans","white" if index==0 else "text_0",True,"center")
    return regions


def render_pptx(manifest: dict[str, Any], output_path: Path) -> list[dict[str, Any]]:
    generator = _load_module("editorial_generate_pptx", GENERATOR_PATH)
    deck = generator.EditorialDeck(style_profile=manifest["deck_profile"], manifest=manifest)
    signatures=[]
    for spec in manifest["slides"]:
        regions=_render_pptx_slide(deck,spec)
        signatures.append({"id":spec["id"],"composition_id":spec["composition_id"],"regions":[list(bounds) for bounds in regions]})
    deck.prs.core_properties.comments=f"manifest_sha256={_manifest_hash(manifest)}"
    output_path.parent.mkdir(parents=True,exist_ok=True); deck.save(output_path)
    return signatures


def render_manifest(manifest: dict[str, Any], output_dir: Path, html_output=True, pptx_output=True) -> dict[str, Any]:
    planner=_load_module("editorial_plan_deck",PLANNER_PATH)
    validation=planner.validate_manifest(manifest)
    if not validation["passed"]: raise ValueError(f"Manifest validation failed: {validation['errors']}")
    output_dir.mkdir(parents=True,exist_ok=True); digest=_manifest_hash(manifest)
    report={"manifest_sha256":digest,"deck_profile":manifest["deck_profile"],"slide_count":len(manifest["slides"]),"outputs":{}}
    if html_output:
        path=output_dir/"deck.html"; report["outputs"]["html"]={"path":str(path),"signatures":render_html(manifest,path)}
    if pptx_output:
        path=output_dir/"deck.pptx"; report["outputs"]["pptx"]={"path":str(path),"signatures":render_pptx(manifest,path)}
    (output_dir/"render-report.json").write_text(json.dumps(report,ensure_ascii=False,indent=2)+"\n",encoding="utf-8")
    return report


def main() -> int:
    parser=argparse.ArgumentParser(description=__doc__); parser.add_argument("manifest",type=Path); parser.add_argument("--output-dir",type=Path,required=True); parser.add_argument("--html-only",action="store_true"); parser.add_argument("--pptx-only",action="store_true"); args=parser.parse_args()
    if args.html_only and args.pptx_only: parser.error("Choose at most one of --html-only/--pptx-only")
    manifest=json.loads(args.manifest.read_text(encoding="utf-8")); report=render_manifest(manifest,args.output_dir,html_output=not args.pptx_only,pptx_output=not args.html_only); print(json.dumps({"manifest_sha256":report["manifest_sha256"],"outputs":list(report["outputs"])},ensure_ascii=False)); return 0


if __name__ == "__main__": raise SystemExit(main())
